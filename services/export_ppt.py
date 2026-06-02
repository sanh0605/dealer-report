import io
from datetime import datetime

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from services.export_excel import (
    REGION_ORDER, PRODUCT_GROUP_LABELS, PRODUCT_GROUP_ORDER,
    _get_ytd_data, _get_month_range, _calc_regional_stats,
)

# --- Colors ---
DARK_BLUE = RGBColor(0x1F, 0x4E, 0x79)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
LIGHT_YELLOW = RGBColor(0xFF, 0xF2, 0xCC)
BLACK = RGBColor(0x00, 0x00, 0x00)
RED = RGBColor(0xFF, 0x00, 0x00)
BLUE_ACCENT = RGBColor(0x25, 0x63, 0xEB)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def _fmt_num(val):
    if isinstance(val, (int, float)):
        if val < 0:
            return f"-{abs(val):,.0f}"
        return f"{val:,.0f}"
    return str(val)


def _fmt_pct(val):
    if isinstance(val, (int, float)):
        return f"{val:.1%}"
    return str(val)


def _add_title(slide, text, top=Inches(0.2)):
    txBox = slide.shapes.add_textbox(Inches(0.4), top, Inches(12), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE


def _add_table(slide, rows, cols, left, top, width, height):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    return table_shape.table


def _set_cell(table, row, col, text, bold=False, color=BLACK, bg=None, size=Pt(10), align=PP_ALIGN.CENTER):
    cell = table.cell(row, col)
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = str(text) if text is not None else ""
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.text = str(text) if text is not None else ""
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    if bg:
        from pptx.oxml.ns import qn
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn('a:solidFill'), {})
        srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % (bg[0], bg[1], bg[2])})
        solidFill.append(srgbClr)
        tcPr.append(solidFill)


def _build_slide_kehoach(prs, ytd_data, reference_date, targets_df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "KẾ HOẠCH 2026")

    dl_data = ytd_data[ytd_data["region"] != "Siêu thị"]
    st_data = ytd_data[ytd_data["region"] == "Siêu thị"]

    rows_data = [
        ("Doanh số", "Đại lý", dl_data["sales_revenue"].sum()),
        ("Doanh số", "Siêu thị", st_data["sales_revenue"].sum()),
        ("Số lượng xe bán ra", "Đại lý", int(dl_data[dl_data["category"].str.upper() == "BIKES"]["sales_volume"].sum())),
        ("Số lượng xe bán ra", "Siêu thị", int(st_data[st_data["category"].str.upper() == "BIKES"]["sales_volume"].sum())),
    ]

    num_rows = len(rows_data) + 1
    tbl = _add_table(slide, num_rows, 5, Inches(0.5), Inches(1.2), Inches(12), Inches(0.4 * num_rows))

    headers = ["Chỉ số", "Kênh áp dụng", "Target", "Kết Quả", "Tiến độ"]
    for i, h in enumerate(headers):
        _set_cell(tbl, 0, i, h, bold=True, color=WHITE, bg=(0x1F, 0x4E, 0x79))

    for r, (chi_so, kenh, ket_qua) in enumerate(rows_data, 1):
        target_val = 0
        if targets_df is not None and not targets_df.empty:
            t = targets_df[(targets_df["kpi"] == chi_so) & (targets_df["channel"] == kenh)]
            if not t.empty:
                target_val = t["target_value"].iloc[0]
        tien_do = ket_qua / target_val if target_val > 0 else 0

        _set_cell(tbl, r, 0, chi_so, align=PP_ALIGN.LEFT)
        _set_cell(tbl, r, 1, kenh)
        _set_cell(tbl, r, 2, _fmt_num(target_val))
        _set_cell(tbl, r, 3, _fmt_num(ket_qua), bold=True, color=BLUE_ACCENT)
        _set_cell(tbl, r, 4, _fmt_pct(tien_do))


def _build_slide_monthly(prs, ytd_data, main_data, reference_date, targets_df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    month_label = f"Tháng {reference_date.month:02d}/{reference_date.year}"
    _add_title(slide, month_label)

    ref = pd.Timestamp(reference_date)
    monthly_data = main_data[
        (main_data["date_transfer"].dt.month == ref.month) &
        (main_data["date_transfer"].dt.year == ref.year)
    ]
    stats_monthly = _calc_regional_stats(monthly_data, targets_df)
    months = _get_month_range(reference_date)

    num_cols = 6 + len(months) + 1  # Khu Vuc + DS + Target + %Target + LNG + %LNG + monthly %LNG + YTD
    active_regions = [rg for rg in REGION_ORDER if rg in stats_monthly]
    num_rows = len(active_regions) + 1

    col_w = Inches(12) / num_cols
    tbl = _add_table(slide, num_rows, num_cols, Inches(0.3), Inches(1.0), Inches(12.3), Inches(0.35 * num_rows))

    headers = ["Khu Vực", "Doanh Số", "Target", "% Target", "LNG", "% LNG"]
    for m_num, _ in months:
        headers.append(f"T{m_num:02d}")
    headers.append("YTD")

    for i, h in enumerate(headers):
        _set_cell(tbl, 0, i, h, bold=True, color=WHITE, bg=(0x1F, 0x4E, 0x79), size=Pt(8))

    for r, rg in enumerate(active_regions, 1):
        sm = stats_monthly[rg]
        is_total = rg == "Grand Total"
        bg = (0xFF, 0xF2, 0xCC) if is_total else None
        txt_color = BLACK

        _set_cell(tbl, r, 0, rg, bold=is_total, bg=bg, align=PP_ALIGN.LEFT, size=Pt(8))
        _set_cell(tbl, r, 1, _fmt_num(sm["revenue"]), bold=is_total,
                  color=RED if sm["revenue"] < 0 else txt_color, bg=bg, size=Pt(8))
        _set_cell(tbl, r, 2, _fmt_num(sm["target"]), bold=is_total, bg=bg, size=Pt(8))
        _set_cell(tbl, r, 3, _fmt_pct(sm["target_pct"]), bold=is_total, bg=bg, size=Pt(8))
        _set_cell(tbl, r, 4, _fmt_num(sm["lng"]), bold=is_total,
                  color=RED if sm["lng"] < 0 else txt_color, bg=bg, size=Pt(8))
        _set_cell(tbl, r, 5, _fmt_pct(sm["lng_pct"]), bold=is_total, bg=bg, size=Pt(8))

        col = 6
        for m_num, _ in months:
            m_data = ytd_data[
                (ytd_data["date_transfer"].dt.month == m_num) &
                (ytd_data["date_transfer"].dt.year == ref.year)
            ]
            m_non_st = m_data[m_data["region"] != "Siêu thị"]
            if rg == "Siêu thị":
                m_data = m_data[m_data["region"] == "Siêu thị"]
            elif rg == "Grand Total":
                pass
            else:
                m_data = m_non_st[m_non_st["region_group"] == rg]
            m_rev = m_data["sales_revenue"].sum()
            m_cost = m_data["cost_of_goods"].sum()
            m_lng_pct = (m_rev - m_cost) / m_rev if m_rev != 0 else 0
            _set_cell(tbl, r, col, _fmt_pct(m_lng_pct), bold=is_total, bg=bg, size=Pt(8))
            col += 1

        stats_ytd = _calc_regional_stats(ytd_data, targets_df)
        sy = stats_ytd.get(rg, {"lng_pct": 0})
        _set_cell(tbl, r, col, _fmt_pct(sy["lng_pct"]), bold=is_total, bg=bg, size=Pt(8))


def _build_slide_luyke(prs, ytd_data, reference_date, targets_df, ar_df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ytd_label = f"Luỹ kế từ 01/01/{reference_date.year} - {reference_date.strftime('%d/%m/%Y')}"
    _add_title(slide, ytd_label)

    ytd_stats = _calc_regional_stats(ytd_data, targets_df)

    if ar_df is not None and not ar_df.empty and "order_id" in ytd_data.columns:
        ar_merged = ar_df.merge(ytd_data[["order_id"]].drop_duplicates(), on="order_id", how="inner")
        ar_merged["doanh_thu"] = ar_merged["Deduction Amout"].fillna(0) + ar_merged["Paid Amout"].fillna(0)
        ar_with_dealer = ar_merged.merge(ytd_data[["order_id", "region", "region_group"]], on="order_id", how="left")
        ar_by_region = ar_with_dealer.groupby("region_group")["doanh_thu"].sum().to_dict()
        ar_sieuthi = ar_with_dealer[ar_with_dealer["region"] == "Siêu thị"]["doanh_thu"].sum()
        ar_total = ar_with_dealer["doanh_thu"].sum()
    else:
        ar_by_region = {}
        ar_sieuthi = 0
        ar_total = 0

    active_regions = [rg for rg in REGION_ORDER if rg in ytd_stats]
    num_rows = len(active_regions) + 1
    tbl = _add_table(slide, num_rows, 6, Inches(0.5), Inches(1.2), Inches(12), Inches(0.4 * num_rows))

    headers = ["Khu Vực", "Doanh số", "Doanh thu", "% Thu hồi", "Target", "% Target"]
    for i, h in enumerate(headers):
        _set_cell(tbl, 0, i, h, bold=True, color=WHITE, bg=(0x1F, 0x4E, 0x79))

    for r, rg in enumerate(active_regions, 1):
        ys = ytd_stats[rg]
        doanh_thu = ar_by_region.get(rg, 0)
        if rg == "Siêu thị":
            doanh_thu = ar_sieuthi
        elif rg == "Grand Total":
            doanh_thu = ar_total
        thu_hoi = doanh_thu / ys["revenue"] if ys["revenue"] != 0 else 0

        is_total = rg == "Grand Total"
        bg = (0xFF, 0xF2, 0xCC) if is_total else None
        _set_cell(tbl, r, 0, rg, bold=is_total, bg=bg, align=PP_ALIGN.LEFT)
        _set_cell(tbl, r, 1, _fmt_num(ys["revenue"]), bold=is_total,
                  color=RED if ys["revenue"] < 0 else BLACK, bg=bg)
        _set_cell(tbl, r, 2, _fmt_num(doanh_thu), bold=is_total,
                  color=RED if doanh_thu < 0 else BLACK, bg=bg)
        _set_cell(tbl, r, 3, _fmt_pct(thu_hoi), bold=is_total, bg=bg)
        _set_cell(tbl, r, 4, _fmt_num(ys["target"]), bold=is_total, bg=bg)
        _set_cell(tbl, r, 5, _fmt_pct(ys["target_pct"]), bold=is_total, bg=bg)


def _build_slide_nhom_san_pham(prs, ytd_data, reference_date):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "NHÓM SẢN PHẨM")

    months = _get_month_range(reference_date)
    num_cols = 1 + len(months) + 1  # Label + months + Grand Total
    num_rows = sum(4 for _ in PRODUCT_GROUP_ORDER) + 1  # 3 rows per group + 1 header
    col_w = Inches(12) / num_cols

    tbl = _add_table(slide, num_rows, num_cols, Inches(0.3), Inches(1.0), Inches(12.3), Inches(0.35 * num_rows))

    headers = ["NHÓM SẢN PHẨM"]
    for m_num, _ in months:
        headers.append(f"T{m_num:02d}")
    headers.append("Grand Total")

    for i, h in enumerate(headers):
        _set_cell(tbl, 0, i, h, bold=True, color=WHITE, bg=(0x1F, 0x4E, 0x79), size=Pt(9))

    r = 1
    for pg_key in PRODUCT_GROUP_ORDER:
        pg_label = PRODUCT_GROUP_LABELS.get(pg_key, pg_key)
        pg_data = ytd_data[ytd_data["product_group"] == pg_key]
        pg_total_rev = pg_data["sales_revenue"].sum()
        pg_total_vol = pg_data["sales_volume"].sum()
        pg_total_cost = pg_data["cost_of_goods"].sum()
        pg_total_lng_pct = (pg_total_rev - pg_total_cost) / pg_total_rev if pg_total_rev != 0 else 0

        _set_cell(tbl, r, 0, pg_label, bold=True, bg=(0xF2, 0xF2, 0xF2), align=PP_ALIGN.LEFT, size=Pt(9))
        for c in range(1, num_cols):
            _set_cell(tbl, r, c, "", bg=(0xF2, 0xF2, 0xF2), size=Pt(9))
        r += 1

        for metric_label, total_val, fmt_fn in [
            ("Doanh số", pg_total_rev, _fmt_num),
            ("Số lượng", pg_total_vol, lambda v: f"{int(v):,}"),
            ("% LNG", pg_total_lng_pct, _fmt_pct),
        ]:
            _set_cell(tbl, r, 0, f"  {metric_label}", align=PP_ALIGN.LEFT, size=Pt(9))
            col = 1
            for m_num, _ in months:
                m_data = pg_data[
                    (pg_data["date_transfer"].dt.month == m_num) &
                    (pg_data["date_transfer"].dt.year == reference_date.year)
                ]
                if metric_label == "Doanh số":
                    val = m_data["sales_revenue"].sum()
                elif metric_label == "Số lượng":
                    val = int(m_data["sales_volume"].sum())
                else:
                    m_rev = m_data["sales_revenue"].sum()
                    m_cost = m_data["cost_of_goods"].sum()
                    val = (m_rev - m_cost) / m_rev if m_rev != 0 else 0
                color = RED if isinstance(val, (int, float)) and val < 0 else BLACK
                _set_cell(tbl, r, col, fmt_fn(val), color=color, size=Pt(9))
                col += 1
            color = RED if isinstance(total_val, (int, float)) and total_val < 0 else BLACK
            _set_cell(tbl, r, col, fmt_fn(total_val), color=color, size=Pt(9))
            r += 1


def _build_top10_overall_slide(prs, ytd_data, reference_date):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "TOP 10 ĐẠI LÝ THEO DOANH SỐ")

    months = _get_month_range(reference_date)
    num_cols = 2 + len(months)  # Name + months + Grand Total

    dealer_monthly = {}
    for _, row in ytd_data.iterrows():
        dname = row.get("dealer_name", "Unknown")
        m = row["date_transfer"].month
        if dname not in dealer_monthly:
            dealer_monthly[dname] = {"monthly": {}, "total": 0}
        dealer_monthly[dname]["monthly"][m] = dealer_monthly[dname]["monthly"].get(m, 0) + row["sales_revenue"]
        dealer_monthly[dname]["total"] += row["sales_revenue"]

    top10 = sorted(dealer_monthly.items(), key=lambda x: x[1]["total"], reverse=True)[:10]
    num_rows = len(top10) + 2  # header + top10 + grand total

    tbl = _add_table(slide, num_rows, num_cols, Inches(0.3), Inches(1.0), Inches(12.3), Inches(0.35 * num_rows))

    headers = ["ĐẠI LÝ"]
    for m_num, _ in months:
        headers.append(f"T{m_num:02d}")
    headers.append("Grand Total")

    for i, h in enumerate(headers):
        _set_cell(tbl, 0, i, h, bold=True, color=WHITE, bg=(0x1F, 0x4E, 0x79), size=Pt(9))

    gt_monthly = {m: 0 for m, _ in months}
    gt_total = 0

    for r, (dname, ddata) in enumerate(top10, 1):
        _set_cell(tbl, r, 0, dname, align=PP_ALIGN.LEFT, size=Pt(9))
        col = 1
        for m_num, _ in months:
            val = ddata["monthly"].get(m_num, 0)
            _set_cell(tbl, r, col, _fmt_num(val), color=RED if val < 0 else BLACK, size=Pt(9))
            gt_monthly[(m_num,)] = gt_monthly.get((m_num,), 0) + val
            col += 1
        _set_cell(tbl, r, col, _fmt_num(ddata["total"]),
                  bold=True, color=BLUE_ACCENT, size=Pt(9))
        gt_total += ddata["total"]

    # Grand total row
    r = num_rows - 1
    _set_cell(tbl, r, 0, "Grand Total", bold=True, bg=(0xFF, 0xF2, 0xCC), align=PP_ALIGN.LEFT, size=Pt(9))
    col = 1
    for m_num, _ in months:
        val = gt_monthly.get((m_num,), 0)
        _set_cell(tbl, r, col, _fmt_num(val), bold=True, bg=(0xFF, 0xF2, 0xCC), size=Pt(9))
        col += 1
    _set_cell(tbl, r, col, _fmt_num(gt_total), bold=True, bg=(0xFF, 0xF2, 0xCC), color=BLUE_ACCENT, size=Pt(9))


def _build_top10_brand_slide(prs, ytd_data, reference_date, product_group, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, title)

    months = _get_month_range(reference_date)
    data = ytd_data[ytd_data["product_group"] == product_group]

    dealer_monthly = {}
    for _, row in data.iterrows():
        dname = row.get("dealer_name", "Unknown")
        m = row["date_transfer"].month
        if dname not in dealer_monthly:
            dealer_monthly[dname] = {"monthly": {}, "total_rev": 0, "total_vol": 0}
        dealer_monthly[dname]["monthly"][m] = dealer_monthly[dname]["monthly"].get(m, {"rev": 0, "vol": 0})
        dealer_monthly[dname]["monthly"][m]["rev"] += row["sales_revenue"]
        dealer_monthly[dname]["monthly"][m]["vol"] += row["sales_volume"]
        dealer_monthly[dname]["total_rev"] += row["sales_revenue"]
        dealer_monthly[dname]["total_vol"] += row["sales_volume"]

    top10 = sorted(dealer_monthly.items(), key=lambda x: x[1]["total_rev"], reverse=True)[:10]
    num_rows = len(top10) * 3 + 2  # each dealer: name row + DS row + SL row + header + GT
    num_cols = 2 + len(months)

    tbl = _add_table(slide, num_rows, num_cols, Inches(0.3), Inches(1.0), Inches(12.3), Inches(0.3 * num_rows))

    headers = ["ĐẠI LÝ"]
    for m_num, _ in months:
        headers.append(f"T{m_num:02d}")
    headers.append("Grand Total")

    for i, h in enumerate(headers):
        _set_cell(tbl, 0, i, h, bold=True, color=WHITE, bg=(0x1F, 0x4E, 0x79), size=Pt(9))

    r = 1
    for dname, ddata in top10:
        _set_cell(tbl, r, 0, dname, bold=True, bg=(0xF2, 0xF2, 0xF2), align=PP_ALIGN.LEFT, size=Pt(9))
        for c in range(1, num_cols):
            _set_cell(tbl, r, c, "", bg=(0xF2, 0xF2, 0xF2), size=Pt(9))
        r += 1

        for metric, total_val in [("rev", ddata["total_rev"]), ("vol", ddata["total_vol"])]:
            label = "Doanh số" if metric == "rev" else "Số lượng"
            fmt_fn = _fmt_num if metric == "rev" else lambda v: f"{int(v):,}"
            _set_cell(tbl, r, 0, f"  {label}", align=PP_ALIGN.LEFT, size=Pt(9))
            col = 1
            for m_num, _ in months:
                val = ddata["monthly"].get(m_num, {}).get(metric, 0)
                color = RED if isinstance(val, (int, float)) and val < 0 else BLACK
                _set_cell(tbl, r, col, fmt_fn(val), color=color, size=Pt(9))
                col += 1
            color = RED if isinstance(total_val, (int, float)) and total_val < 0 else BLACK
            _set_cell(tbl, r, col, fmt_fn(total_val), bold=True, color=color, size=Pt(9))
            r += 1

    # Grand total row
    _set_cell(tbl, r, 0, "Grand Total", bold=True, bg=(0xFF, 0xF2, 0xCC), align=PP_ALIGN.LEFT, size=Pt(9))
    col = 1
    for m_num, _ in months:
        _set_cell(tbl, r, col, "", bold=True, bg=(0xFF, 0xF2, 0xCC), size=Pt(9))
        col += 1
    _set_cell(tbl, r, col, "", bold=True, bg=(0xFF, 0xF2, 0xCC), size=Pt(9))


def generate_ppt_bytes(
    filtered_data: pd.DataFrame,
    previous_filtered_data: pd.DataFrame,
    main_data: pd.DataFrame,
    reference_date: datetime,
) -> bytes:
    from database.gsheets_db import read_sheet

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    targets_df = read_sheet("sales_targets")
    ar_df = read_sheet("accounts_receivable")

    ref = pd.Timestamp(reference_date)
    ytd_data = _get_ytd_data(main_data, ref)

    # Slide 1: Ke hoach 2026
    _build_slide_kehoach(prs, ytd_data, ref, targets_df)

    # Slide 2: Monthly regional
    _build_slide_monthly(prs, ytd_data, main_data, ref, targets_df)

    # Slide 3: Luy ke
    _build_slide_luyke(prs, ytd_data, ref, targets_df, ar_df)

    # Slide 4: Nhom san pham
    _build_slide_nhom_san_pham(prs, ytd_data, ref)

    # Slide 5: Top 10 overall
    _build_top10_overall_slide(prs, ytd_data, ref)

    # Slides 6-9: Top 10 per brand
    brand_slides = [
        ("Xe OEM", "TOP 10 ĐẠI LÝ THEO DOANH SỐ XE OEM"),
        ("Xe Java", "TOP 10 ĐẠI LÝ THEO DOANH SỐ XE JAVA"),
        ("Xe Giant", "TOP 10 ĐẠI LÝ THEO DOANH SỐ XE GIANT"),
        ("Maxxis", "TOP 10 ĐẠI LÝ THEO DOANH SỐ MAXXIS"),
    ]
    for pg_key, title in brand_slides:
        _build_top10_brand_slide(prs, ytd_data, ref, product_group=pg_key, title=title)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
